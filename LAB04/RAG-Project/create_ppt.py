import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_box(slide, left, top, width, height, text, subtext, color=RGBColor(59, 130, 246)):
    # Add a rectangle
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, 
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(250, 250, 250)
    shape.line.color.rgb = color
    shape.line.width = Pt(2)
    
    # Add text
    tf = shape.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(12)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    
    if subtext:
        p2 = tf.add_paragraph()
        p2.text = subtext
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(100, 100, 100)
        p2.alignment = PP_ALIGN.CENTER
        
    return shape

def add_arrow(slide, left, top, width=0.4):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, 
        Inches(left), Inches(top), Inches(width), Inches(0.15)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(150, 150, 150)
    shape.line.color.rgb = RGBColor(150, 150, 150)
    return shape

def add_section_label(slide, left, top, text):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(2), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(50, 50, 50)


def create_presentation():
    prs = Presentation()
    # Use 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ----------------------------------------------------
    # SLIDE 1: Workflow Diagram
    # ----------------------------------------------------
    blank_slide_layout = prs.slide_layouts[5] # Title Only
    slide1 = prs.slides.add_slide(blank_slide_layout)
    shapes = slide1.shapes
    
    title = slide1.shapes.title
    title.text = "RAG System Workflow Diagram"
    
    # Colors
    c_blue = RGBColor(59, 130, 246)
    c_green = RGBColor(16, 185, 129)
    c_orange = RGBColor(249, 115, 22)
    c_purple = RGBColor(139, 92, 246)

    # --- SECTION 1: BUILD ---
    add_section_label(slide1, 0.5, 1.2, "1. BUILD - Prepare dataset (python build_index.py)")
    
    y = 1.7
    add_box(slide1, 0.5, y, 1.8, 0.8, "Knowledge Base", "data/cat_qa_dataset.txt", c_blue)
    add_arrow(slide1, 2.4, y+0.35)
    add_box(slide1, 2.9, y, 1.8, 0.8, "document_loader", "Parse Q&A", c_blue)
    add_arrow(slide1, 4.8, y+0.35)
    add_box(slide1, 5.3, y, 1.8, 0.8, "text_splitter", "Chunk size 400", c_blue)
    add_arrow(slide1, 7.2, y+0.35)
    add_box(slide1, 7.7, y, 1.8, 0.8, "embedding_model", "MiniLM-L12-v2", c_blue)
    add_arrow(slide1, 9.6, y+0.35)
    add_box(slide1, 10.1, y, 1.8, 0.8, "vector_store", "FAISS / BM25", c_blue)
    
    # --- SECTION 2: QUERY ---
    add_section_label(slide1, 0.5, 3.2, "2. QUERY - Answer questions (python main.py)")
    
    y = 3.8
    # Input -> Retrieval -> Context -> LLM -> Output
    add_box(slide1, 0.5, y, 1.7, 1.0, "User Question", "Input via terminal", c_green)
    add_arrow(slide1, 2.3, y+0.45)
    
    add_box(slide1, 2.8, y, 1.7, 1.0, "query_transform", "Processing (Optional)", c_green)
    add_arrow(slide1, 4.6, y+0.45)
    
    add_box(slide1, 5.1, y, 1.7, 1.0, "hybrid_retriever", "Retrieval\nBM25 + Dense FAISS", c_orange)
    add_arrow(slide1, 6.9, y+0.45)
    
    add_box(slide1, 7.4, y, 1.7, 1.0, "rerankers", "Refinement\nbge-reranker-v2", c_green)
    add_arrow(slide1, 9.2, y+0.45)
    
    add_box(slide1, 9.7, y, 1.7, 1.0, "generator", "Context + LLM\nollama / openai", c_orange)
    add_arrow(slide1, 11.5, y+0.45, width=0.25)
    
    add_box(slide1, 11.85, y, 1.3, 1.0, "Answer", "Final Output", c_green)
    
    # Memory block
    add_box(slide1, 6.25, y+1.3, 2.0, 0.6, "memory", "Stores last 6 turns", c_orange)
    # Adding a simple text for memory flow
    txBox = slide1.shapes.add_textbox(Inches(4.5), Inches(y+1.4), Inches(2), Inches(0.5))
    txBox.text_frame.text = "<-- Feedback loop for context"
    txBox.text_frame.paragraphs[0].font.size = Pt(10)

    # --- SECTION 3: EVALUATION ---
    add_section_label(slide1, 0.5, 5.8, "3. EVALUATION - Measure and improve")
    
    y = 6.3
    add_box(slide1, 0.5, y, 1.8, 0.8, "golden_set.json", "Ground Truth Data", c_purple)
    add_arrow(slide1, 2.4, y+0.35)
    add_box(slide1, 2.9, y, 1.8, 0.8, "eval_retrieval.py", "Hit@k, MRR, nDCG", c_purple)
    add_arrow(slide1, 4.8, y+0.35)
    add_box(slide1, 5.3, y, 1.8, 0.8, "eval_generation.py", "Faithfulness, Correctness", c_purple)
    add_arrow(slide1, 7.2, y+0.35)
    add_box(slide1, 7.7, y, 1.8, 0.8, "metrics.py", "Compare Configs", c_purple)


    # ----------------------------------------------------
    # SLIDE 2: Workflow Explanation
    # ----------------------------------------------------
    bullet_slide_layout = prs.slide_layouts[1] # Title and Content
    slide2 = prs.slides.add_slide(bullet_slide_layout)
    
    title = slide2.shapes.title
    title.text = "RAG Workflow Explanation"
    
    body_shape = slide2.shapes.placeholders[1]
    tf = body_shape.text_frame
    tf.text = "1. Input (User Question & query_transform)"
    
    p = tf.add_paragraph()
    p.text = "User types a question in the terminal (main.py)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Question is optionally transformed (rewritten, multi-query) before searching."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. Retrieval (hybrid_retriever & rerankers)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Searches the FAISS Vector Database (Semantic) and BM25 Index (Keyword)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Combines results using Reciprocal Rank Fusion (RRF) and optionally reranks."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. Context (generator)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "Top chunks are combined into a numbered reference block (Context)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Inserted into the Prompt Template along with conversation history."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "4. LLM (generator)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "The prompt + context is sent to the LLM (ollama, openai, or gemini)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "LLM generates an answer strictly based on the reference data."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "5. Output (Answer)"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "System appends medical disclaimers and lists [n] citations to the generated text."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Final answer is displayed to the user."
    p.level = 1

    # Format font size for bullet points to fit nicely
    for paragraph in tf.paragraphs:
        if paragraph.level == 0:
            paragraph.font.size = Pt(20)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(59, 130, 246)
        else:
            paragraph.font.size = Pt(16)

    prs.save('D:\\ATCS-CPE\\LAB03\\RAG-Project\\RAG_Workflow_Presentation.pptx')

if __name__ == '__main__':
    create_presentation()
